from pptx import Presentation
from pptx.util import Inches, Pt
from src.engine.layout import LayoutDeck

def render_pptx(layout_deck: LayoutDeck, filename: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for slide_data in layout_deck.slides:
        slide = prs.slides.add_slide(blank_layout)
        
        # Notes
        if slide_data.notes:
            slide.notes_slide.notes_text_frame.text = slide_data.notes

        for box in slide_data.boxes:
            if box.kind == "text":
                txBox = slide.shapes.add_textbox(Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
                tf = txBox.text_frame
                tf.word_wrap = True
                if box.text:
                    p = tf.paragraphs[0]
                    p.text = box.text
                    p.font.size = Pt(box.font_size)
            
            elif box.kind == "image" and box.image_ref and box.image_ref.local_path:
                try:
                    slide.shapes.add_picture(box.image_ref.local_path, Inches(box.x), Inches(box.y), width=Inches(box.w), height=Inches(box.h))
                except Exception as e:
                    print(f"⚠️ Erro ao inserir imagem {box.image_ref.local_path}: {e}")

    prs.save(filename)
    return filename